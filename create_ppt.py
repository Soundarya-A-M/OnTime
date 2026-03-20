"""
OnTime - Real-Time Bus Tracking System
Professional PowerPoint Presentation Generator (12 Slides)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──────────────────────────────────────────────
DARK_BG      = RGBColor(0x0F, 0x17, 0x2A)
ACCENT_BLUE  = RGBColor(0x38, 0xBD, 0xF8)
ACCENT_PURPLE= RGBColor(0xA7, 0x8B, 0xFA)
ACCENT_GREEN = RGBColor(0x4A, 0xDE, 0x80)
ACCENT_ORANGE= RGBColor(0xFB, 0xBF, 0x24)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCB, 0xD5, 0xE1)
MID_GRAY     = RGBColor(0x94, 0xA3, 0xB8)
CARD_BG      = RGBColor(0x1E, 0x29, 0x3B)
RED          = RGBColor(0xF8, 0x71, 0x71)
CYAN         = RGBColor(0x22, 0xD3, 0xEE)
PINK         = RGBColor(0xF4, 0x72, 0xB6)

TOTAL_SLIDES = 12

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


# ── Helpers ────────────────────────────────────────────────────

def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, fill_color=CARD_BG, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_text(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return p

def add_paragraph(tf, text, size=16, color=LIGHT_GRAY, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(4), font_name="Segoe UI"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p

def add_title_bar(slide, title_text, subtitle_text=None):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.5), Inches(0.15), Inches(0.7))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()
    tb = add_text_box(slide, Inches(1.15), Inches(0.4), Inches(10), Inches(0.8))
    set_text(tb.text_frame, title_text, size=32, color=WHITE, bold=True)
    if subtitle_text:
        tb2 = add_text_box(slide, Inches(1.15), Inches(1.1), Inches(10), Inches(0.5))
        set_text(tb2.text_frame, subtitle_text, size=16, color=MID_GRAY)
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Pt(2))
    sep.fill.solid()
    sep.fill.fore_color.rgb = RGBColor(0x33, 0x44, 0x55)
    sep.line.fill.background()

def add_bullet_card(slide, left, top, width, height, title, bullets, accent=ACCENT_BLUE):
    card = add_shape_rect(slide, left, top, width, height, CARD_BG, border_color=RGBColor(0x33, 0x44, 0x55))
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()
    tb = add_text_box(slide, left + Inches(0.25), top + Inches(0.25), width - Inches(0.5), height - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, title, size=18, color=accent, bold=True)
    for b in bullets:
        add_paragraph(tf, f"•  {b}", size=14, color=LIGHT_GRAY)

def slide_number_footer(slide, num, total=TOTAL_SLIDES):
    tb = add_text_box(slide, Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4))
    set_text(tb.text_frame, f"{num} / {total}", size=11, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

def draw_er_box(slide, left, top, width, height, title, fields, accent):
    """Draw a mini entity box for the ER diagram."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1A, 0x25, 0x36)
    card.line.color.rgb = accent
    card.line.width = Pt(1.5)
    # Header bar
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.35))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = accent
    hdr.line.fill.background()
    tb = add_text_box(slide, left + Inches(0.08), top + Inches(0.03), width - Inches(0.15), Inches(0.3))
    set_text(tb.text_frame, title, size=12, color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)
    # Fields
    tb2 = add_text_box(slide, left + Inches(0.08), top + Inches(0.38), width - Inches(0.15), height - Inches(0.45))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    first = True
    for f in fields:
        if first:
            set_text(tf2, f, size=9, color=LIGHT_GRAY)
            first = False
        else:
            add_paragraph(tf2, f, size=9, color=LIGHT_GRAY, space_before=Pt(1), space_after=Pt(1))

def draw_arch_box(slide, left, top, width, height, label, accent):
    """Draw a rounded architecture layer box."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x1A, 0x25, 0x36)
    box.line.color.rgb = accent
    box.line.width = Pt(1.5)
    box.text_frame.paragraphs[0].text = label
    box.text_frame.paragraphs[0].font.size = Pt(11)
    box.text_frame.paragraphs[0].font.color.rgb = accent
    box.text_frame.paragraphs[0].font.bold = True
    box.text_frame.paragraphs[0].font.name = "Segoe UI"
    box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def draw_arrow_down(slide, cx, top, length, color=MID_GRAY):
    """Draw a simple vertical connector line."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, top, Pt(2), length)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    # Arrow head (small triangle approximation)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, cx - Pt(5), top + length, Pt(12), Pt(10))
    tri.fill.solid()
    tri.fill.fore_color.rgb = color
    tri.line.fill.background()


# ══════════════════════════════════════════════════════════════
#  SLIDE 1 – Title Slide
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

tb = add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.2))
tf = tb.text_frame
set_text(tf, "OnTime", size=60, color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "Real-Time Bus Tracking System", size=28, color=WHITE, bold=False, space_before=Pt(8))

tb2 = add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.5))
set_text(tb2.text_frame, "GPS Tracking  •  Route Visualization  •  ETA Prediction  •  Seat Booking", size=16, color=MID_GRAY)

card = add_shape_rect(slide, Inches(1), Inches(4.0), Inches(5.5), Inches(2.8), CARD_BG, RGBColor(0x33, 0x44, 0x55))
tb3 = add_text_box(slide, Inches(1.3), Inches(4.15), Inches(5), Inches(2.6))
tf3 = tb3.text_frame; tf3.word_wrap = True
set_text(tf3, "Submitted By", size=14, color=ACCENT_BLUE, bold=True)
add_paragraph(tf3, "1.  Santhosh Kumar R S", size=16, color=WHITE)
add_paragraph(tf3, "2.  [Student Name 2]", size=16, color=WHITE)
add_paragraph(tf3, "3.  [Student Name 3]", size=16, color=WHITE)
add_paragraph(tf3, "4.  [Student Name 4]", size=16, color=WHITE)

card2 = add_shape_rect(slide, Inches(7), Inches(4.0), Inches(5.5), Inches(2.8), CARD_BG, RGBColor(0x33, 0x44, 0x55))
tb4 = add_text_box(slide, Inches(7.3), Inches(4.15), Inches(5), Inches(2.6))
tf4 = tb4.text_frame; tf4.word_wrap = True
set_text(tf4, "Submitted To", size=14, color=ACCENT_PURPLE, bold=True)
add_paragraph(tf4, "[Guide / Professor Name]", size=16, color=WHITE, bold=True)
add_paragraph(tf4, "[Designation]", size=14, color=MID_GRAY)
add_paragraph(tf4, "[Department / College Name]", size=14, color=MID_GRAY)

slide_number_footer(slide, 1)


# ══════════════════════════════════════════════════════════════
#  SLIDE 2 – Abstract
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "Abstract")

abstract = (
    "OnTime is a comprehensive full-stack web application designed to revolutionize urban public "
    "transportation through real-time bus tracking. Built on the MERN stack (MongoDB, Express.js, "
    "React.js, Node.js), the system leverages WebSocket technology via Socket.IO to deliver live "
    "GPS location updates every 3–5 seconds, enabling passengers to track buses on an interactive "
    "Mapbox-powered map.\n\n"
    "The platform features an intelligent ETA prediction engine that calculates arrival times based "
    "on real-time speed and distance metrics. A complete seat booking module allows passengers to "
    "select seats, generate digital tickets, and manage their booking history.\n\n"
    "Role-based dashboards provide tailored experiences for three user types: Passengers can track "
    "buses and book tickets; Drivers can start/end trips with automatic GPS sharing; and Admins can "
    "manage buses, routes, driver assignments, and monitor system analytics. Security is ensured "
    "through JWT authentication, bcrypt password hashing, CORS protection, and API rate limiting.\n\n"
    "The application embodies a modern smart-city aesthetic with glassmorphism effects, dark-mode "
    "gradients, and a fully responsive design optimized for both desktop and mobile devices."
)

tb = add_text_box(slide, Inches(1), Inches(1.9), Inches(11.3), Inches(5.0))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, abstract, size=16, color=LIGHT_GRAY)
slide_number_footer(slide, 2)


# ══════════════════════════════════════════════════════════════
#  SLIDE 3 – Problem Statement (with Example Scenario)
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "Problem Statement")

# ── Example scenario card (left) ──
example_card = add_shape_rect(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(5.0), CARD_BG, RGBColor(0x33, 0x44, 0x55))
strip_e = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.9), Inches(5.5), Pt(4))
strip_e.fill.solid(); strip_e.fill.fore_color.rgb = ACCENT_ORANGE; strip_e.line.fill.background()

tb_e = add_text_box(slide, Inches(1.05), Inches(2.1), Inches(5.0), Inches(4.6))
tf_e = tb_e.text_frame; tf_e.word_wrap = True
set_text(tf_e, "Real-World Example", size=18, color=ACCENT_ORANGE, bold=True)
add_paragraph(tf_e, "", size=6, color=CARD_BG, space_before=Pt(4), space_after=Pt(2))

scenario_lines = [
    ("Scenario:", WHITE, True),
    ("Priya, a college student, leaves her home every morning to catch Bus No. 42 from the Koyambedu bus stop.", LIGHT_GRAY, False),
    ("", CARD_BG, False),
    ("The Problem:", ACCENT_ORANGE, True),
    ("•  She has no idea whether the bus has already passed or is delayed.", LIGHT_GRAY, False),
    ("•  She waits 30+ minutes in the sun with no information.", LIGHT_GRAY, False),
    ("•  She can't book a seat — the bus arrives overcrowded.", LIGHT_GRAY, False),
    ("•  Missing the bus means she's late to college — again.", LIGHT_GRAY, False),
    ("", CARD_BG, False),
    ("With OnTime:", ACCENT_GREEN, True),
    ("She opens the app, sees Bus 42 is 3 stops away (ETA 6 min), books Seat 14, and arrives on time.", ACCENT_GREEN, False),
]

for text, clr, bld in scenario_lines:
    add_paragraph(tf_e, text, size=13, color=clr, bold=bld, space_before=Pt(3), space_after=Pt(2))

# ── Problem Statement card (right) ──
ps_card = add_shape_rect(slide, Inches(6.8), Inches(1.9), Inches(5.5), Inches(5.0), CARD_BG, RGBColor(0x33, 0x44, 0x55))
strip_p = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.9), Inches(5.5), Pt(4))
strip_p.fill.solid(); strip_p.fill.fore_color.rgb = RED; strip_p.line.fill.background()

tb_p = add_text_box(slide, Inches(7.05), Inches(2.1), Inches(5.0), Inches(4.6))
tf_p = tb_p.text_frame; tf_p.word_wrap = True
set_text(tf_p, "Problem Statement", size=18, color=RED, bold=True)
add_paragraph(tf_p, "", size=6, color=CARD_BG, space_before=Pt(4), space_after=Pt(2))

problem_lines = [
    "Urban public bus transit systems in India lack real-time tracking infrastructure, forcing millions of commuters to depend on static, unreliable timetables.",
    "Passengers have no visibility into live bus locations, leading to prolonged and unpredictable waiting times at bus stops.",
    "There is no mechanism for seat reservation, resulting in overcrowded buses and poor commuter experience.",
    "Transit authorities lack a centralized digital platform to manage fleets, monitor drivers, and analyze ridership data.",
    "Existing solutions are either city-locked, proprietary, or do not offer an integrated tracking + booking system.",
    "There is a clear need for a scalable, open, web-based platform that provides real-time GPS tracking, ETA prediction, digital seat booking, and role-based dashboards for all stakeholders.",
]

for pl in problem_lines:
    p = tf_p.add_paragraph()
    p.space_before = Pt(6); p.space_after = Pt(3)
    run_icon = p.add_run(); run_icon.text = "▸  "; run_icon.font.size = Pt(13); run_icon.font.color.rgb = RED; run_icon.font.bold = True; run_icon.font.name = "Segoe UI"
    run_text = p.add_run(); run_text.text = pl; run_text.font.size = Pt(12); run_text.font.color.rgb = LIGHT_GRAY; run_text.font.name = "Segoe UI"

slide_number_footer(slide, 3)


# ══════════════════════════════════════════════════════════════
#  SLIDE 4 – Existing System / Literature Review
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "Existing System / Literature Review")

systems = [
    ("Google Maps Transit", [
        "Shows static schedule-based transit data",
        "Limited to cities with GTFS data feeds",
        "No seat booking or driver interaction",
        "Dependent on third-party transit agencies"
    ]),
    ("City-Specific Apps (BMTC, DTC)", [
        "Tied to a single city's transport authority",
        "Often lack real-time GPS tracking",
        "No role-based management dashboards",
        "Fragmented user experience"
    ]),
    ("SMS / IVRS-Based Systems", [
        "Users must manually request bus status",
        "No visual map or route visualization",
        "High latency in information delivery",
        "No booking capabilities"
    ]),
]

for i, (title, bullets) in enumerate(systems):
    left = Inches(0.8) + i * Inches(4.1)
    add_bullet_card(slide, left, Inches(1.9), Inches(3.8), Inches(4.5), title, bullets,
                    accent=[ACCENT_BLUE, ACCENT_PURPLE, ACCENT_GREEN][i])

slide_number_footer(slide, 4)


# ══════════════════════════════════════════════════════════════
#  SLIDE 5 – Disadvantages of Existing System
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "Disadvantages of Existing System")

disadvantages = [
    "No Real-Time Tracking — Passengers rely on fixed timetables with no live location data.",
    "Inaccurate ETAs — Schedules don't account for traffic, breakdowns, or route deviations.",
    "No Seat Booking — Passengers cannot reserve seats, leading to overcrowding and uncertainty.",
    "Lack of Role-Based Dashboards — No unified platform for drivers, admins, and passengers.",
    "Poor Scalability — City-specific solutions cannot be adapted to other regions easily.",
    "No Driver Accountability — No mechanism to monitor driver trips or routes in real time.",
    "Fragmented Systems — Separate apps for tracking, booking, and management create confusion.",
    "Limited Mobile Responsiveness — Many existing systems are not optimized for mobile devices.",
]

tb = add_text_box(slide, Inches(1), Inches(1.9), Inches(11.3), Inches(5.0))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "", size=14, color=LIGHT_GRAY)

for d in disadvantages:
    p = tf.add_paragraph(); p.space_before = Pt(10); p.space_after = Pt(4)
    run_icon = p.add_run(); run_icon.text = "  ✗   "; run_icon.font.size = Pt(16); run_icon.font.color.rgb = RED; run_icon.font.bold = True; run_icon.font.name = "Segoe UI"
    run_text = p.add_run(); run_text.text = d; run_text.font.size = Pt(15); run_text.font.color.rgb = LIGHT_GRAY; run_text.font.name = "Segoe UI"

slide_number_footer(slide, 5)


# ══════════════════════════════════════════════════════════════
#  SLIDE 6 – Proposed System
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "Proposed System — OnTime")

features = [
    ("Real-Time GPS Tracking", "Live bus location updates every 3–5 sec via Socket.IO WebSockets."),
    ("Interactive Map", "Mapbox-powered map with moving bus markers and route polylines."),
    ("ETA Prediction", "Speed & distance-based arrival time calculation updated in real time."),
    ("Seat Booking", "Full digital booking flow — seat selection, confirmation, e-ticket generation."),
    ("Role-Based Dashboards", "Tailored panels for Passengers, Drivers, and Admins."),
    ("Secure Authentication", "JWT tokens, bcrypt hashing, CORS protection, and rate limiting."),
]

colors = [ACCENT_BLUE, ACCENT_PURPLE, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_BLUE, ACCENT_PURPLE]
for i, (title, desc) in enumerate(features):
    col = i % 3; row = i // 3
    left = Inches(0.8) + col * Inches(4.1)
    top  = Inches(1.9) + row * Inches(2.5)
    card = add_shape_rect(slide, left, top, Inches(3.8), Inches(2.2), CARD_BG, RGBColor(0x33, 0x44, 0x55))
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.8), Pt(4))
    strip.fill.solid(); strip.fill.fore_color.rgb = colors[i]; strip.line.fill.background()
    tb = add_text_box(slide, left + Inches(0.25), top + Inches(0.25), Inches(3.3), Inches(1.8))
    tf = tb.text_frame; tf.word_wrap = True
    set_text(tf, title, size=18, color=colors[i], bold=True)
    add_paragraph(tf, desc, size=14, color=LIGHT_GRAY, space_before=Pt(10))

slide_number_footer(slide, 6)


# ══════════════════════════════════════════════════════════════
#  SLIDE 7 – Advantages of Proposed System
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "Advantages of Proposed System")

advantages = [
    "Real-Time Visibility — Passengers know exactly where their bus is at all times.",
    "Accurate ETAs — Dynamic calculation based on live speed and GPS data.",
    "Seamless Booking — Digital seat selection and e-tickets reduce queues and overcrowding.",
    "Unified Platform — Single app for tracking, booking, and fleet management.",
    "Scalable Architecture — Cloud-hosted MERN stack can scale to any city or fleet size.",
    "Driver Monitoring — Admins can track all active trips and driver performance in real time.",
    "Secure & Reliable — JWT auth, bcrypt encryption, rate limiting, and CORS protection.",
    "Mobile-First Design — Fully responsive UI works seamlessly on phones and tablets.",
]

tb = add_text_box(slide, Inches(1), Inches(1.9), Inches(11.3), Inches(5.0))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "", size=14, color=LIGHT_GRAY)

for adv in advantages:
    p = tf.add_paragraph(); p.space_before = Pt(10); p.space_after = Pt(4)
    run_icon = p.add_run(); run_icon.text = "  ✓   "; run_icon.font.size = Pt(16); run_icon.font.color.rgb = ACCENT_GREEN; run_icon.font.bold = True; run_icon.font.name = "Segoe UI"
    run_text = p.add_run(); run_text.text = adv; run_text.font.size = Pt(15); run_text.font.color.rgb = LIGHT_GRAY; run_text.font.name = "Segoe UI"

slide_number_footer(slide, 7)


# ══════════════════════════════════════════════════════════════
#  SLIDE 8 – Models / Methodology
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "System Model & Methodology")

add_bullet_card(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(5.0),
    "Agile Development Methodology", [
        "Iterative development with sprint cycles",
        "Continuous integration & testing",
        "Modular MVC architecture (Model-View-Controller)",
        "RESTful API design principles",
        "WebSocket event-driven communication",
        "Role-Based Access Control (RBAC)",
    ], accent=ACCENT_BLUE)

add_bullet_card(slide, Inches(6.8), Inches(1.9), Inches(5.5), Inches(5.0),
    "3-Tier Architecture", [
        "Presentation Layer — React.js + Tailwind CSS",
        "Business Logic Layer — Node.js + Express.js",
        "Data Layer — MongoDB Atlas (NoSQL)",
        "Real-Time Layer — Socket.IO (WebSocket)",
        "Auth Layer — JWT + bcryptjs",
        "Map Layer — Mapbox GL / Leaflet",
    ], accent=ACCENT_PURPLE)

slide_number_footer(slide, 8)


# ══════════════════════════════════════════════════════════════
#  SLIDE 9 – Model Explanation (Data Flow)
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "Model Explanation — Data Flow & Modules")

modules = [
    ("User Module", "Registration, login, JWT auth, profile management, role assignment (Passenger / Driver / Admin)."),
    ("Route Module", "CRUD for routes with stops, distances, durations. Admin-managed."),
    ("Bus Module", "Bus registration, capacity management, driver assignment, route linking."),
    ("Trip Module", "Driver starts/ends trips. GPS coordinates broadcast via Socket.IO every 3–5 sec."),
    ("Booking Module", "Passengers select seats, confirm bookings, generate digital e-tickets."),
    ("Real-Time Module", "Socket.IO handles location broadcast, ETA updates, trip status events."),
]

colors2 = [ACCENT_BLUE, ACCENT_PURPLE, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_BLUE, ACCENT_PURPLE]
for i, (title, desc) in enumerate(modules):
    col = i % 3; row = i // 3
    left = Inches(0.8) + col * Inches(4.1)
    top  = Inches(1.9) + row * Inches(2.5)
    card = add_shape_rect(slide, left, top, Inches(3.8), Inches(2.2), CARD_BG, RGBColor(0x33, 0x44, 0x55))
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.2), top + Inches(0.25), Inches(0.45), Inches(0.45))
    badge.fill.solid(); badge.fill.fore_color.rgb = colors2[i]; badge.line.fill.background()
    badge.text_frame.paragraphs[0].text = str(i + 1)
    badge.text_frame.paragraphs[0].font.size = Pt(16)
    badge.text_frame.paragraphs[0].font.color.rgb = DARK_BG
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].font.name = "Segoe UI"
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb = add_text_box(slide, left + Inches(0.8), top + Inches(0.2), Inches(2.8), Inches(1.9))
    tf = tb.text_frame; tf.word_wrap = True
    set_text(tf, title, size=17, color=colors2[i], bold=True)
    add_paragraph(tf, desc, size=13, color=LIGHT_GRAY, space_before=Pt(8))

slide_number_footer(slide, 9)


# ══════════════════════════════════════════════════════════════
#  SLIDE 10 – ER Diagram + Architecture (Combined)
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "ER Diagram & System Architecture")

# ── LEFT HALF: ER Diagram ──────────────────────────────────────
# Section label
tb_er = add_text_box(slide, Inches(0.8), Inches(1.75), Inches(5.8), Inches(0.4))
set_text(tb_er.text_frame, "Entity-Relationship Diagram", size=15, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# Entity boxes  (positioned in a diamond-ish layout)
draw_er_box(slide,
    Inches(2.2), Inches(2.2), Inches(2.2), Inches(1.4),
    "USER", ["_id (PK)", "name, email, password", "phone, role (enum)", "isActive, bookingHistory[]"],
    ACCENT_BLUE)

draw_er_box(slide,
    Inches(0.4), Inches(4.0), Inches(2.2), Inches(1.3),
    "ROUTE", ["_id (PK)", "routeName, routeNumber", "stops[{name, lat, lng}]", "distance, estimatedDuration"],
    ACCENT_PURPLE)

draw_er_box(slide,
    Inches(4.2), Inches(4.0), Inches(2.2), Inches(1.3),
    "BUS", ["_id (PK), busNumber", "routeId (FK→Route)", "driverId (FK→User)", "capacity, status, isOnTrip"],
    ACCENT_GREEN)

draw_er_box(slide,
    Inches(0.4), Inches(5.7), Inches(2.2), Inches(1.3),
    "TRIP", ["_id (PK)", "busId (FK), routeId (FK)", "driverId (FK→User)", "status, bookedSeats[]"],
    ACCENT_ORANGE)

draw_er_box(slide,
    Inches(4.2), Inches(5.7), Inches(2.2), Inches(1.3),
    "BOOKING", ["_id (PK), ticketId", "userId, busId, tripId (FK)", "seatNumber, fromStop", "amount, status"],
    PINK)

# Relationship connectors (simple lines)
# USER → BUS  (driverId)
draw_arrow_down(slide, Inches(4.9), Inches(3.6), Inches(0.4), ACCENT_GREEN)
# USER → TRIP  (driverId)
draw_arrow_down(slide, Inches(1.4), Inches(3.6), Inches(0.4), ACCENT_ORANGE)
# ROUTE → BUS
line_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.6), Inches(4.65), Inches(1.6), Pt(2))
line_h.fill.solid(); line_h.fill.fore_color.rgb = MID_GRAY; line_h.line.fill.background()
# BUS → BOOKING
draw_arrow_down(slide, Inches(5.3), Inches(5.3), Inches(0.4), PINK)
# ROUTE → TRIP
draw_arrow_down(slide, Inches(1.5), Inches(5.3), Inches(0.4), ACCENT_ORANGE)

# ── RIGHT HALF: System Architecture ───────────────────────────
tb_arch = add_text_box(slide, Inches(6.8), Inches(1.75), Inches(5.6), Inches(0.4))
set_text(tb_arch.text_frame, "System Architecture", size=15, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

# Architecture layers (stacked vertically)
arch_x = Inches(7.3)
arch_w = Inches(4.7)
layer_h = Inches(0.65)
gap = Inches(0.15)

layers = [
    ("CLIENT LAYER\nReact.js  •  Tailwind CSS  •  Zustand  •  Leaflet", ACCENT_BLUE),
    ("REAL-TIME LAYER\nSocket.IO Client  ↔  Socket.IO Server", CYAN),
    ("API LAYER\nExpress.js REST API  •  JWT Auth Middleware", ACCENT_GREEN),
    ("BUSINESS LOGIC LAYER\nControllers  •  Validators  •  Rate Limiter", ACCENT_ORANGE),
    ("DATA LAYER\nMongoose ODM  •  MongoDB Atlas (Cloud NoSQL)", ACCENT_PURPLE),
]

for i, (label, accent) in enumerate(layers):
    top = Inches(2.3) + i * (layer_h + gap + Inches(0.25))
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, arch_x, top, arch_w, layer_h + Inches(0.1))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x1A, 0x25, 0x36)
    box.line.color.rgb = accent
    box.line.width = Pt(1.5)
    
    tb_l = add_text_box(slide, arch_x + Inches(0.15), top + Inches(0.05), arch_w - Inches(0.3), layer_h + Inches(0.05))
    tf_l = tb_l.text_frame; tf_l.word_wrap = True
    lines = label.split("\n")
    set_text(tf_l, lines[0], size=11, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    if len(lines) > 1:
        add_paragraph(tf_l, lines[1], size=9, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(2), space_after=Pt(0))

    # Arrow between layers (except last)
    if i < len(layers) - 1:
        arrow_top = top + layer_h + Inches(0.1)
        cx = arch_x + arch_w / 2
        draw_arrow_down(slide, cx, arrow_top, Inches(0.15), MID_GRAY)

slide_number_footer(slide, 10)


# ══════════════════════════════════════════════════════════════
#  SLIDE 11 – Hardware & Software Requirements
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "Hardware & Software Requirements")

add_bullet_card(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(5.0),
    "Hardware Requirements", [
        "Processor: Intel Core i3 or higher",
        "RAM: 4 GB minimum (8 GB recommended)",
        "Storage: 500 MB free disk space",
        "Network: Stable internet connection",
        "GPS: Device with GPS capability (for drivers)",
        "Display: 1366×768 resolution or higher",
    ], accent=ACCENT_BLUE)

add_bullet_card(slide, Inches(6.8), Inches(1.9), Inches(5.5), Inches(5.0),
    "Software Requirements", [
        "OS: Windows 10+ / macOS / Linux",
        "Runtime: Node.js v16 or higher",
        "Database: MongoDB Atlas (Cloud)",
        "Browser: Chrome / Firefox / Edge (latest)",
        "IDE: VS Code (recommended)",
        "Tools: Git, npm, Postman (testing)",
        "Maps API: Mapbox / Leaflet",
        "Deployment: Vercel (frontend), Railway (backend)",
    ], accent=ACCENT_PURPLE)

slide_number_footer(slide, 11)


# ══════════════════════════════════════════════════════════════
#  SLIDE 12 – Conclusion
# ══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "Conclusion")

conclusion = (
    "OnTime is a modern, full-stack web application that addresses the critical challenges "
    "of urban public transportation by providing real-time bus tracking, intelligent ETA "
    "predictions, and a complete digital booking system.\n\n"
    "The application demonstrates the effective use of the MERN stack combined with WebSocket "
    "technology to deliver a seamless, real-time user experience. The role-based architecture "
    "ensures that passengers, drivers, and administrators each have a purpose-built interface "
    "tailored to their needs.\n\n"
    "With its scalable cloud-hosted architecture, robust security measures, and mobile-responsive "
    "design, OnTime can be deployed in any city to improve public transit efficiency, reduce "
    "passenger wait times, and provide transit authorities with powerful fleet management tools."
)

tb = add_text_box(slide, Inches(1), Inches(1.9), Inches(11.3), Inches(3.5))
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, conclusion, size=16, color=LIGHT_GRAY)

add_bullet_card(slide, Inches(1), Inches(5.2), Inches(11.3), Inches(1.8),
    "Future Scope", [
        "Push notifications for delays  •  Payment gateway integration  •  QR code ticket scanning  •  Route optimization algorithms  •  Mobile app (React Native)  •  Multi-language support",
    ], accent=ACCENT_GREEN)

tb = add_text_box(slide, Inches(4.5), Inches(7.0), Inches(5), Inches(0.4))
set_text(tb.text_frame, "Thank You!", size=18, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
slide_number_footer(slide, 12)


# ── Save ──────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(__file__), "OnTime_Presentation_v2.pptx")
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"📊 Total slides: {TOTAL_SLIDES}")
